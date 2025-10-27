#!/usr/bin/env python3
"""
CV Automation Script
====================
This script automates the creation of PowerPoint slides from CV data, instructions,
and optionally a job description.

Features:
- Processes CV documents (PDF, DOCX, or text)
- Reads Excel instruction files for customization
- Uses PowerPoint templates for consistent formatting
- Optionally incorporates job description requirements
- Generates client-facing PowerPoint slides
- Creates traceability reports

Usage:
    python cv_automation.py --cv CV_FILE --instructions EXCEL_FILE --template PPT_TEMPLATE [--job-description JD_FILE]
"""

import argparse
import logging
import sys
from pathlib import Path
from typing import Optional, Dict, Any, List
import json
from datetime import datetime

# Third-party imports (to be installed via requirements.txt)
try:
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import openpyxl
except ImportError as e:
    print(f"Missing required package: {e}")
    print("Please install required packages: pip install -r requirements.txt")
    sys.exit(1)


class CVAutomation:
    """Main class for CV automation workflow."""
    
    def __init__(self, cv_path: Path, instructions_path: Path, 
                 template_path: Path, job_description_path: Optional[Path] = None,
                 output_dir: Path = None):
        """
        Initialize the CV automation system.
        
        Args:
            cv_path: Path to CV file (PDF, DOCX, or TXT)
            instructions_path: Path to Excel instruction file
            template_path: Path to PowerPoint template
            job_description_path: Optional path to job description file
            output_dir: Directory for output files (default: current directory)
        """
        self.cv_path = cv_path
        self.instructions_path = instructions_path
        self.template_path = template_path
        self.job_description_path = job_description_path
        self.output_dir = output_dir or Path.cwd()
        
        # Initialize data structures
        self.cv_data: Dict[str, Any] = {}
        self.instructions: Dict[str, Any] = {}
        self.job_requirements: Optional[Dict[str, Any]] = None
        self.traceability_log: List[Dict[str, Any]] = []
        
        # Setup logging
        self._setup_logging()
        
    def _setup_logging(self):
        """Configure logging for the application."""
        log_format = '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
        logging.basicConfig(
            level=logging.INFO,
            format=log_format,
            handlers=[
                logging.StreamHandler(sys.stdout),
                logging.FileHandler(self.output_dir / 'cv_automation.log')
            ]
        )
        self.logger = logging.getLogger(__name__)
        
    def validate_inputs(self) -> bool:
        """
        Validate that all required input files exist and are accessible.
        
        Returns:
            True if all validations pass, False otherwise
        """
        self.logger.info("Validating input files...")
        
        # Check CV file
        if not self.cv_path.exists():
            self.logger.error(f"CV file not found: {self.cv_path}")
            return False
            
        # Check instructions file
        if not self.instructions_path.exists():
            self.logger.error(f"Instructions file not found: {self.instructions_path}")
            return False
        
        if not self.instructions_path.suffix.lower() in ['.xlsx', '.xls']:
            self.logger.error(f"Instructions file must be Excel format: {self.instructions_path}")
            return False
            
        # Check template file
        if not self.template_path.exists():
            self.logger.error(f"Template file not found: {self.template_path}")
            return False
            
        if not self.template_path.suffix.lower() == '.pptx':
            self.logger.error(f"Template must be PowerPoint format: {self.template_path}")
            return False
            
        # Check job description (optional)
        if self.job_description_path and not self.job_description_path.exists():
            self.logger.error(f"Job description file not found: {self.job_description_path}")
            return False
            
        self.logger.info("All input files validated successfully")
        return True
        
    def parse_cv(self) -> Dict[str, Any]:
        """
        Parse CV file and extract structured information.
        
        Returns:
            Dictionary containing parsed CV data
        """
        self.logger.info(f"Parsing CV: {self.cv_path}")
        
        cv_data = {
            'file_path': str(self.cv_path),
            'file_name': self.cv_path.name,
            'parsed_date': datetime.now().isoformat(),
            'sections': {},
            'raw_text': ''
        }
        
        # Parse based on file type
        file_extension = self.cv_path.suffix.lower()
        
        if file_extension == '.txt':
            cv_data['raw_text'] = self.cv_path.read_text(encoding='utf-8')
        elif file_extension in ['.pdf', '.docx']:
            # Placeholder for PDF/DOCX parsing
            # Would use libraries like PyPDF2, pdfplumber, or python-docx
            self.logger.warning(f"PDF/DOCX parsing not yet implemented. Using filename only.")
            cv_data['raw_text'] = f"CV from file: {self.cv_path.name}"
        else:
            self.logger.warning(f"Unsupported CV format: {file_extension}")
            cv_data['raw_text'] = f"Unsupported format: {self.cv_path.name}"
            
        # Extract basic sections (simplified for now)
        cv_data['sections'] = self._extract_cv_sections(cv_data['raw_text'])
        
        self.cv_data = cv_data
        self._log_traceability('cv_parsing', 'Parsed CV file', cv_data)
        
        return cv_data
        
    def _extract_cv_sections(self, text: str) -> Dict[str, str]:
        """
        Extract sections from CV text.
        
        Args:
            text: Raw CV text
            
        Returns:
            Dictionary of section names to content
        """
        # Basic section extraction (to be enhanced)
        sections = {
            'summary': '',
            'experience': '',
            'education': '',
            'skills': '',
            'certifications': ''
        }
        
        # Placeholder implementation
        # In a real implementation, would use NLP or regex patterns
        sections['summary'] = text[:200] if len(text) > 200 else text
        
        return sections
        
    def parse_instructions(self) -> Dict[str, Any]:
        """
        Parse Excel instruction file.
        
        Returns:
            Dictionary containing instruction data
        """
        self.logger.info(f"Parsing instructions: {self.instructions_path}")
        
        try:
            # Read Excel file
            df = pd.read_excel(self.instructions_path)
            
            instructions = {
                'file_path': str(self.instructions_path),
                'parsed_date': datetime.now().isoformat(),
                'rules': [],
                'mappings': {},
                'raw_data': df.to_dict('records')
            }
            
            # Process instructions based on expected format
            # Assuming columns: Section, Field, Rule, Value
            for _, row in df.iterrows():
                rule = {
                    'section': row.get('Section', ''),
                    'field': row.get('Field', ''),
                    'rule': row.get('Rule', ''),
                    'value': row.get('Value', '')
                }
                instructions['rules'].append(rule)
                
            self.instructions = instructions
            self._log_traceability('instruction_parsing', 'Parsed instruction file', instructions)
            
            return instructions
            
        except Exception as e:
            self.logger.error(f"Error parsing instructions: {e}")
            raise
            
    def parse_job_description(self) -> Optional[Dict[str, Any]]:
        """
        Parse job description file if provided.
        
        Returns:
            Dictionary containing job description data or None
        """
        if not self.job_description_path:
            self.logger.info("No job description provided")
            return None
            
        self.logger.info(f"Parsing job description: {self.job_description_path}")
        
        try:
            jd_text = self.job_description_path.read_text(encoding='utf-8')
            
            job_requirements = {
                'file_path': str(self.job_description_path),
                'parsed_date': datetime.now().isoformat(),
                'raw_text': jd_text,
                'requirements': [],
                'key_skills': []
            }
            
            # Basic extraction (to be enhanced)
            job_requirements['requirements'] = [
                line.strip() for line in jd_text.split('\n') 
                if line.strip() and len(line.strip()) > 10
            ]
            
            self.job_requirements = job_requirements
            self._log_traceability('jd_parsing', 'Parsed job description', job_requirements)
            
            return job_requirements
            
        except Exception as e:
            self.logger.error(f"Error parsing job description: {e}")
            raise
            
    def generate_powerpoint(self, output_path: Path = None) -> Path:
        """
        Generate PowerPoint presentation from CV data and instructions.
        
        Args:
            output_path: Optional custom output path
            
        Returns:
            Path to generated PowerPoint file
        """
        self.logger.info("Generating PowerPoint presentation")
        
        if output_path is None:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = self.output_dir / f"CV_Presentation_{timestamp}.pptx"
            
        try:
            # Load template
            prs = Presentation(str(self.template_path))
            
            # Add or modify slides based on instructions
            self._populate_slides(prs)
            
            # Save presentation
            prs.save(str(output_path))
            
            self.logger.info(f"PowerPoint generated: {output_path}")
            self._log_traceability('ppt_generation', 'Generated PowerPoint', {'output_path': str(output_path)})
            
            return output_path
            
        except Exception as e:
            self.logger.error(f"Error generating PowerPoint: {e}")
            raise
            
    def _populate_slides(self, prs: Presentation):
        """
        Populate PowerPoint slides with CV data according to instructions.
        
        Args:
            prs: PowerPoint Presentation object
        """
        # Get or create first slide
        if len(prs.slides) == 0:
            # Add a blank slide if template has no slides
            # Find a suitable blank layout (usually last, but varies by template)
            blank_layout_index = min(6, len(prs.slide_layouts) - 1)
            blank_layout = prs.slide_layouts[blank_layout_index]
            slide = prs.slides.add_slide(blank_layout)
        else:
            slide = prs.slides[0]
            
        # Add title if not present
        if not slide.shapes.title:
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "CV Summary"
            
        # Add CV content based on instructions
        # This is a simplified implementation
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Add summary from CV
        if self.cv_data and 'sections' in self.cv_data:
            summary = self.cv_data['sections'].get('summary', 'No summary available')
            tf.text = f"Summary:\n{summary}"
            
    def generate_traceability_report(self, output_path: Path = None) -> Path:
        """
        Generate traceability report documenting all processing steps.
        
        Args:
            output_path: Optional custom output path
            
        Returns:
            Path to generated report file
        """
        self.logger.info("Generating traceability report")
        
        if output_path is None:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = self.output_dir / f"Traceability_Report_{timestamp}.json"
            
        report = {
            'generated_at': datetime.now().isoformat(),
            'inputs': {
                'cv': str(self.cv_path),
                'instructions': str(self.instructions_path),
                'template': str(self.template_path),
                'job_description': str(self.job_description_path) if self.job_description_path else None
            },
            'processing_log': self.traceability_log,
            'summary': {
                'total_steps': len(self.traceability_log),
                'cv_sections_extracted': len(self.cv_data.get('sections', {})),
                'instruction_rules': len(self.instructions.get('rules', []))
            }
        }
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
            
        self.logger.info(f"Traceability report generated: {output_path}")
        
        return output_path
        
    def _log_traceability(self, step: str, description: str, data: Any):
        """
        Log a traceability entry.
        
        Args:
            step: Step identifier
            description: Human-readable description
            data: Associated data (will be summarized for logging)
        """
        # Create a meaningful summary of the data
        data_str = str(data)
        max_length = 500  # Configurable truncation length
        
        if len(data_str) > max_length:
            # Truncate with ellipsis to indicate there's more
            data_summary = data_str[:max_length] + '...'
        else:
            data_summary = data_str
            
        entry = {
            'timestamp': datetime.now().isoformat(),
            'step': step,
            'description': description,
            'data_summary': data_summary
        }
        self.traceability_log.append(entry)
        
    def run(self) -> Dict[str, Path]:
        """
        Execute the complete CV automation workflow.
        
        Returns:
            Dictionary containing paths to generated files
        """
        self.logger.info("="*60)
        self.logger.info("Starting CV Automation Workflow")
        self.logger.info("="*60)
        
        # Validate inputs
        if not self.validate_inputs():
            raise ValueError("Input validation failed")
            
        # Parse all inputs
        self.parse_cv()
        self.parse_instructions()
        self.parse_job_description()
        
        # Generate outputs
        ppt_path = self.generate_powerpoint()
        report_path = self.generate_traceability_report()
        
        results = {
            'powerpoint': ppt_path,
            'traceability_report': report_path
        }
        
        self.logger.info("="*60)
        self.logger.info("CV Automation Workflow Completed Successfully")
        self.logger.info(f"PowerPoint: {ppt_path}")
        self.logger.info(f"Report: {report_path}")
        self.logger.info("="*60)
        
        return results


def main():
    """Main entry point for the script."""
    parser = argparse.ArgumentParser(
        description='Automate PowerPoint slide creation from CV data',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
Examples:
  # Basic usage
  python cv_automation.py --cv resume.pdf --instructions rules.xlsx --template template.pptx
  
  # With job description
  python cv_automation.py --cv resume.pdf --instructions rules.xlsx --template template.pptx --job-description jd.txt
  
  # With custom output directory
  python cv_automation.py --cv resume.pdf --instructions rules.xlsx --template template.pptx --output-dir ./output
        '''
    )
    
    parser.add_argument('--cv', required=True, type=Path,
                        help='Path to CV file (PDF, DOCX, or TXT)')
    parser.add_argument('--instructions', required=True, type=Path,
                        help='Path to Excel instruction file')
    parser.add_argument('--template', required=True, type=Path,
                        help='Path to PowerPoint template file')
    parser.add_argument('--job-description', type=Path,
                        help='Optional path to job description file')
    parser.add_argument('--output-dir', type=Path, default=Path.cwd(),
                        help='Output directory for generated files (default: current directory)')
    parser.add_argument('--verbose', '-v', action='store_true',
                        help='Enable verbose logging')
    
    args = parser.parse_args()
    
    # Set logging level
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
        
    # Create output directory if it doesn't exist
    args.output_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        # Initialize and run automation
        automation = CVAutomation(
            cv_path=args.cv,
            instructions_path=args.instructions,
            template_path=args.template,
            job_description_path=args.job_description,
            output_dir=args.output_dir
        )
        
        results = automation.run()
        
        print("\n" + "="*60)
        print("SUCCESS: CV Automation completed!")
        print("="*60)
        print(f"PowerPoint: {results['powerpoint']}")
        print(f"Traceability Report: {results['traceability_report']}")
        print("="*60 + "\n")
        
        return 0
        
    except Exception as e:
        print(f"\nERROR: {e}", file=sys.stderr)
        logging.error(f"Fatal error: {e}", exc_info=True)
        return 1


if __name__ == '__main__':
    sys.exit(main())
